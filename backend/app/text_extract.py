"""Extract plain text from PDF, DOCX, PPTX, XLSX (best-effort)."""

from __future__ import annotations

import io
from typing import BinaryIO

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


def extract_text(data: bytes, suffix: str) -> str:
    s = suffix.lower()
    bio: BinaryIO = io.BytesIO(data)
    if s == ".pdf":
        return _pdf(bio)
    if s == ".docx":
        return _docx(bio)
    if s == ".pptx":
        return _pptx(bio)
    if s == ".xlsx":
        return _xlsx(bio)
    return ""


def _pdf(bio: BinaryIO) -> str:
    reader = PdfReader(bio)
    parts: list[str] = []
    for page in reader.pages:
        t = page.extract_text() or ""
        if t.strip():
            parts.append(t)
    return "\n\n".join(parts).strip()


def _docx(bio: BinaryIO) -> str:
    doc = Document(bio)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()


def _pptx(bio: BinaryIO) -> str:
    prs = Presentation(bio)
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines).strip()


def _xlsx(bio: BinaryIO) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ""
    wb = load_workbook(bio, read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append("\t".join(cells))
    return "\n".join(lines).strip()
